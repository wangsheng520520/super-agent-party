import json
import os
import re
import sys
from urllib.parse import urlparse
import aiohttp
from io import BytesIO
import asyncio
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook
from striprtf.striprtf import rtf_to_text
from odf import text
from odf.opendocument import load  # ODF 处理移动到这里避免重复导入
from pptx import Presentation
from urllib.parse import urlparse, urlunparse
from py.get_setting  import get_host,get_port
import zipfile
import xml.etree.ElementTree as ET
# 平台检测
IS_WINDOWS = sys.platform == 'win32'
IS_MAC = sys.platform == 'darwin'

# 动态文件类型配置
BASE_OFFICE_EXTS = ['doc', 'docx', 'pptx', 'xls', 'xlsx', 'pdf', 'rtf', 'odt', 'epub']
PLATFORM_SPECIFIC_EXTS = {
    'win32': ['ppt'],
    'darwin': ['pages', 'numbers', 'key']
}

FILE_FILTERS = [
    { 
        'name': '办公文档', 
        'extensions': BASE_OFFICE_EXTS + PLATFORM_SPECIFIC_EXTS.get(sys.platform, [])
    },
    { 
        'name': '编程开发', 
        'extensions': [
            'js', 'ts', 'py', 'java', 'c', 'cpp', 'h', 'hpp', 'go', 'rs',
            'swift', 'kt', 'dart', 'rb', 'php', 'html', 'css', 'scss',
            'less', 'vue', 'svelte', 'jsx', 'tsx', 'json', 'xml', 'yml',
            'yaml', 'sql', 'sh'
        ]
    },
    {
        'name': '数据配置',
        'extensions': ['csv', 'tsv', 'txt', 'md', 'log', 'conf', 'ini', 'env', 'toml']
    }
]

office_extensions = {ext for group in FILE_FILTERS if group['name'] == '办公文档' for ext in group['extensions']}

async def handle_url(url):
    """异步处理URL输入，优先尝试替换 host:port 的内部URL，失败后回退到原始URL"""
    parsed_url = urlparse(url)
    original_url = url  # 保存原始URL用于回退
    modified_url = url  # 默认使用原URL

    # 如果路径包含 'uploaded_files'，则替换 host:port
    if 'uploaded_files' in parsed_url.path:
        HOST = get_host()
        PORT = get_port()
        if HOST == '0.0.0.0':
            HOST = '127.0.0.1'
        modified_parsed_url = parsed_url._replace(netloc=f'{HOST}:{PORT}')
        modified_url = urlunparse(modified_parsed_url)

    # 尝试使用修改后的URL请求
    async with aiohttp.ClientSession() as session:
        for try_url in [modified_url, original_url]:  # 最多两次尝试：先改写URL，再原始URL
            try:
                async with session.get(try_url, headers={
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3'
                }) as response:
                    response.raise_for_status()
                    content = await response.read()
                    ext = os.path.splitext(parsed_url.path)[1].lstrip('.').lower()
                    return content, ext
            except (aiohttp.ClientError, OSError) as e:
                print(f"请求失败（{try_url}）: {e}")
                if try_url == original_url:
                    raise  # 如果是最后一次尝试失败，抛出异常
                else:
                    continue  # 否则继续尝试原始URL

async def handle_local_file(file_path):
    """异步处理本地文件"""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    loop = asyncio.get_event_loop()
    content = await loop.run_in_executor(None, _read_file, file_path)
    ext = os.path.splitext(file_path)[1].lstrip('.').lower()
    return content, ext

def _read_file(file_path):
    """同步读取文件内容"""
    with open(file_path, 'rb') as f:
        return f.read()

async def get_content(input_str):
    """获取文件内容和扩展名"""
    if input_str.startswith(('http://', 'https://')):
        return await handle_url(input_str)
    else:
        return await handle_local_file(input_str)

def decode_text(content_bytes):
    """通用文本解码（增加BOM处理）"""
    encodings = ['utf-8-sig', 'utf-16', 'gbk', 'iso-8859-1', 'latin-1']
    for enc in encodings:
        try:
            return content_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return content_bytes.decode('utf-8', errors='replace')

async def handle_office_document(content, ext):
    """异步处理办公文档（带平台检测）"""
    handler = {
        'pdf': handle_pdf,
        'docx': handle_docx,
        'xlsx': handle_excel,
        'xls': handle_excel,
        'rtf': handle_rtf,
        'odt': handle_odt,
        'pptx': handle_pptx,
        'epub': handle_epub,  # 添加epub处理
    }
    
    # Windows平台扩展
    if IS_WINDOWS:
        handler['ppt'] = handle_ppt
        handler['doc'] = handle_doc
    
    handler_func = handler.get(ext)
    
    if handler_func:
        return await handler_func(content)
    
    # Mac平台iWork格式处理
    if IS_MAC and ext in ['pages', 'numbers', 'key']:
        raise NotImplementedError(f"iWork格式暂不支持自动解析，请手动导出为通用格式")
    
    raise NotImplementedError(f"暂不支持处理 {ext.upper()} 格式文件")

# 添加EPUB处理函数
async def handle_epub(content):
    """异步处理EPUB文件"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_epub, content)

import posixpath  # 新增导入

def _process_epub(content):
    """同步处理EPUB内容，返回JSON格式的章节结构"""
    try:
        chapters = []
        processed_files = set()  # 用于记录已处理的文件路径

        with BytesIO(content) as epub_file:
            with zipfile.ZipFile(epub_file, 'r') as epub_zip:
                # 解析容器文件获取OPF路径
                container_data = epub_zip.read('META-INF/container.xml')
                container_root = ET.fromstring(container_data)
                opf_path_element = container_root.find('.//{*}rootfile')
                if opf_path_element is None:
                    raise ValueError("OPF文件路径未找到")
                opf_path = opf_path_element.get('full-path')

                # 解析OPF文件
                opf_data = epub_zip.read(opf_path)
                opf_root = ET.fromstring(opf_data)
                opf_namespace = {'opf': 'http://www.idpf.org/2007/opf'}
                
                # 获取spine顺序（章节阅读顺序）
                spine = opf_root.find('.//opf:spine', opf_namespace)
                if spine is None:
                    raise ValueError("spine元素未找到")
                itemrefs = [item.get('idref') for item in spine.findall('opf:itemref', opf_namespace)]
                
                # 构建manifest映射 (id -> file路径)
                manifest = {}
                for item in opf_root.findall('.//opf:item', opf_namespace):
                    item_id = item.get('id')
                    href = item.get('href')
                    if item_id and href:
                        # 使用posixpath处理路径
                        manifest[item_id] = posixpath.normpath(href)
                
                # OPF文件所在目录
                opf_dir = posixpath.dirname(opf_path)
                
                # 按spine顺序处理每个章节
                for item_id in itemrefs:
                    if item_id not in manifest:
                        continue
                    
                    # 使用posixpath拼接路径
                    rel_path = manifest[item_id]
                    abs_path = posixpath.join(opf_dir, rel_path) if opf_dir else rel_path
                    abs_path = posixpath.normpath(abs_path)

                    # 查找实际存在的文件名（解决大小写敏感问题）
                    actual_path = None
                    for name in epub_zip.namelist():
                        if name.replace('\\', '/').lower() == abs_path.lower().replace('\\', '/'):
                            actual_path = name
                            break
                    
                    # 如果文件已处理过，跳过
                    if actual_path in processed_files:
                        continue
                    
                    if actual_path and actual_path in epub_zip.namelist():
                        with epub_zip.open(actual_path) as chapter_file:
                            html_data = chapter_file.read()
                            chapter_title, chapter_text = _parse_epub_chapter(html_data)
                            chapter_content = f"{chapter_title}\n\n{chapter_text}" if chapter_title else chapter_text
                            if chapter_content.strip():
                                chapters.append(chapter_content)
                            processed_files.add(actual_path)  # 标记为已处理
        
        return json.dumps({"chapters": chapters}, ensure_ascii=False)
    
    except Exception as e:
        raise RuntimeError(f"EPUB解析失败: {str(e)}")



def _parse_epub_chapter(html_data):
    """解析单个章节内容，返回(标题, 正文)"""
    try:
        root = ET.fromstring(html_data)
        ns = {'xhtml': 'http://www.w3.org/1999/xhtml'}
        
        # 1. 提取标题
        title = ""
        for level in range(1, 7):
            title_elem = root.find(f'.//xhtml:h{level}', ns)
            if title_elem is not None and title_elem.text:
                title = title_elem.text.strip()
                found_level = level  # 记录实际找到的标题级别
                break
        else:
            found_level = 0  # 未找到标题
        
        # 2. 提取正文（精确控制提取范围）
        body_text = []
        
        # 方案一：直接提取整个 body 内容（推荐）
        body_elem = root.find('.//xhtml:body', ns)
        if body_elem is not None:
            # 提取所有文本（自动合并子元素）
            full_text = ''.join(body_elem.itertext()).strip()
            if full_text:
                body_text.append(full_text)
        
        # 3. 过滤标题内容（如果标题在 body 中）
        final_text = []
        for text in body_text:
            # 移除标题行（如果有）
            cleaned = text.replace(title, '', 1).strip()
            final_text.append(cleaned if cleaned else text)
        
        return title, '\n'.join(final_text).strip()

    except ET.ParseError:
        # 备选方案：正则表达式处理
        html_str = html_data.decode('utf-8', errors='replace')
        title_match = re.search(r'<h[1-6][^>]*>(.*?)</h[1-6]>', html_str, re.IGNORECASE)
        title = title_match.group(1).strip() if title_match else ""
        
        # 提取 body 内容
        body_match = re.search(r'<body[^>]*>(.*?)</body>', html_str, re.DOTALL | re.IGNORECASE)
        body_content = body_match.group(1) if body_match else html_str
        
        # 去除所有标签
        text = re.sub(r'<[^>]+>', '', body_content).strip()
        return title, text


def _extract_text_from_xml_element(element):
    """递归提取XML元素中的文本"""
    text_parts = []
    
    # 添加元素的文本内容
    if element.text and element.text.strip():
        text_parts.append(element.text.strip())
    
    # 递归处理子元素
    for child in element:
        text_parts.append(_extract_text_from_xml_element(child))
    
    # 添加元素的尾部文本
    if element.tail and element.tail.strip():
        text_parts.append(element.tail.strip())
    
    return ' '.join(text_parts)


async def handle_odt(content):
    """异步处理ODT文件"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_odt, content)

def _process_odt(content):
    """同步处理ODT内容"""
    from odf.teletype import extractText
    
    try:
        doc = load(BytesIO(content))
        text_content = []
        for para in doc.getElementsByType(text.P):
            text_content.append(extractText(para))
        for table in doc.getElementsByType(text.Table):
            for row in table.getElementsByType(text.TableRow):
                row_data = []
                for cell in row.getElementsByType(text.TableCell):
                    row_data.append(extractText(cell))
                text_content.append("\t".join(row_data))
        return '\n'.join(text_content)
    except Exception as e:
        raise RuntimeError(f"ODT文件解析失败: {str(e)}")

async def handle_pdf(content):
    """异步处理PDF文件（增加容错处理）"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_pdf, content)

def _process_pdf(content):
    """同步处理PDF内容"""
    text = []
    try:
        with BytesIO(content) as pdf_file:
            reader = PdfReader(pdf_file)
            for page in reader.pages:
                page_text = page.extract_text() or ""  # 处理无文本页面
                text.append(page_text)
    except Exception as e:
        raise RuntimeError(f"PDF解析失败: {str(e)}")
    return '\n'.join(text)

async def handle_docx(content):
    """异步处理DOCX文件"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_docx, content)

def _process_docx(content):
    """同步处理DOCX内容（增加表格处理）"""
    doc = Document(BytesIO(content))
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            text.append('\t'.join(cell.text for cell in row.cells))
    return '\n'.join(text)

async def handle_excel(content):
    """异步处理Excel文件（优化大文件处理）"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_excel, content)

def _process_excel(content):
    """同步处理Excel内容"""
    try:
        wb = load_workbook(filename=BytesIO(content), read_only=True, data_only=True)
        text = []
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                text.append('\t'.join(str(cell) if cell is not None else '' for cell in row))
        return '\n'.join(text)
    except Exception as e:
        raise RuntimeError(f"Excel解析失败: {str(e)}")

async def handle_rtf(content):
    """异步处理RTF文件"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_rtf, content)

def _process_rtf(content):
    """同步处理RTF内容"""
    try:
        return rtf_to_text(content.decode('utf-8', errors='replace'))
    except Exception as e:
        raise RuntimeError(f"RTF解析失败: {str(e)}")

async def handle_pptx(content):
    """异步处理PPTX文件（优化内容提取）"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_pptx, content)

def _process_pptx(content):
    """同步处理PPTX内容"""
    try:
        prs = Presentation(BytesIO(content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        row_data = [cell.text_frame.text.strip() for cell in row.cells]
                        text.append("\t".join(row_data))
        return '\n'.join(filter(None, text))
    except Exception as e:
        raise RuntimeError(f"PPTX解析失败: {str(e)}")

async def handle_ppt(content):
    """处理PPT文件（Windows平台专用）"""
    if not IS_WINDOWS:
        raise NotImplementedError("PPT格式仅支持在Windows系统处理")
    
    try:
        import win32com.client
    except ImportError:
        raise RuntimeError("请安装pywin32依赖: pip install pywin32")
    
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_ppt, content)

def _process_ppt(content):
    """同步处理PPT内容（Windows COM API）"""
    import win32com.client
    import tempfile
    import pythoncom

    pythoncom.CoInitialize()
    try:
        with tempfile.NamedTemporaryFile(suffix='.ppt', delete=False) as tmp_file:
            tmp_file.write(content)
            tmp_path = tmp_file.name
        
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        pres = powerpoint.Presentations.Open(tmp_path)
        text = []
        for slide in pres.Slides:
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    text.append(shape.TextFrame.TextRange.Text.strip())
        pres.Close()
        powerpoint.Quit()
        return '\n'.join(filter(None, text))
    except Exception as e:
        raise RuntimeError(f"PPT解析失败: {str(e)}")
    finally:
        pythoncom.CoUninitialize()
        os.unlink(tmp_path)

# 2. 实现 handle_doc 函数
async def handle_doc(content):
    if not IS_WINDOWS:
        raise NotImplementedError("DOC格式仅支持在Windows系统处理")
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_doc, content)

def _process_doc(content):
    import win32com.client
    import tempfile
    import pythoncom
    
    pythoncom.CoInitialize()
    try:
        with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as tmp_file:
            tmp_file.write(content)
            tmp_path = tmp_file.name
            
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(tmp_path)
        text = doc.Range().Text
        doc.Close()
        word.Quit()
        return text.strip()
    except Exception as e:
        raise RuntimeError(f"DOC解析失败: {str(e)}")
    finally:
        pythoncom.CoUninitialize()
        if 'tmp_path' in locals():
            os.unlink(tmp_path)

async def get_file_content(file_url):
    """异步获取文件内容（增加编码异常处理）"""
    try:
        content, ext = await get_content(file_url)
        if ext in office_extensions:
            return await handle_office_document(content, ext)
        return decode_text(content)
    except Exception as e:
        return f"文件解析错误: {str(e)}"

async def get_files_content(files_path_list):
    """异步获取所有文件内容并拼接（增加错误隔离）"""
    tasks = [get_file_content(fp) for fp in files_path_list]
    contents = await asyncio.gather(*tasks, return_exceptions=True)
    results = []
    for fp, content in zip(files_path_list, contents):
        if isinstance(content, Exception):
            results.append(f"文件 {fp} 解析失败: {str(content)}")
        else:
            results.append(f"文件 {fp} 内容：\n{content}")
    return "\n\n".join(results)

async def get_files_json(files_list):
    """异步获取所有文件内容并拼接为JSON格式（增加错误隔离）
    输入
    files_list: [{'path': 'path/to/file', 'name': 'file_name'}]
    """
    tasks = [get_file_content(files["path"]) for files in files_list]
    contents = await asyncio.gather(*tasks, return_exceptions=True)
    results = []
    for files, content in zip(files_list, contents):
        results.append({"file_path": files["path"],"file_name": files["name"], "content": str(content)})
    return results

ALLOWED_EXTENSIONS = [
  # 办公文档
    'doc', 'docx', 'ppt', 'pptx', 'xls', 'xlsx', 'pdf', 'pages', 
    'numbers', 'key', 'rtf', 'odt', 'epub',
  
  # 编程开发
  'js', 'ts', 'py', 'java', 'c', 'cpp', 'h', 'hpp', 'go', 'rs',
  'swift', 'kt', 'dart', 'rb', 'php', 'html', 'css', 'scss', 'less',
  'vue', 'svelte', 'jsx', 'tsx', 'json', 'xml', 'yml', 'yaml', 
  'sql', 'sh',
  
  # 数据配置
  'csv', 'tsv', 'txt', 'md', 'log', 'conf', 'ini', 'env', 'toml'
]

ALLOWED_IMAGE_EXTENSIONS = ['png', 'jpg', 'jpeg', 'gif', 'webp', 'bmp']

file_tool = {
    "type": "function",
    "function": {
        "name": "get_file_content",
        "description": f"获取给定的文件URL中的内容，无论是公网URL还是服务器内部URL，由于工具调用结果会被缓存在服务器中，本工具也可以通过工具调用结果的URL用来查看工具调用结果，支持格式：{', '.join(ALLOWED_EXTENSIONS)}",
        "parameters": {
            "type": "object",
            "properties": {
                "file_url": {
                    "type": "string",
                    "description": "文件URL或者工具调用结果的URL",
                }
            },
            "required": ["file_url"],
        },
    },
}

image_tool = {
    "type": "function",
    "function": {
        "name": "get_image_content",
        "description": f"获取给定的图片URL中的内容，无论是公网URL还是服务器内部URL，支持格式：{', '.join(ALLOWED_IMAGE_EXTENSIONS)}",
        "parameters": {
            "type": "object",
            "properties": {
                "image_url": {
                    "type": "string",
                    "description": "图片URL",
                }
            },
            "required": ["image_url"],
        },
    },
}